from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from typing import Dict, Any, List, Optional, Tuple
import io
import os
import random
import asyncio
import concurrent.futures
import hashlib
import json
import logging

from app.services.image_service import ImageServiceFactory

logger = logging.getLogger(__name__)


class ThemeConfig:
    def __init__(
        self,
        primary_color: Tuple[int, int, int] = (0, 112, 192),
        secondary_color: Tuple[int, int, int] = (68, 84, 106),
        background_color: Tuple[int, int, int] = (255, 255, 255),
        title_font: str = "微软雅黑",
        body_font: str = "微软雅黑",
        title_size: int = 44,
        body_size: int = 18,
        accent_color: Tuple[int, int, int] = (255, 192, 0)
    ):
        self.primary_color = RGBColor(*primary_color)
        self.secondary_color = RGBColor(*secondary_color)
        self.background_color = RGBColor(*background_color)
        self.title_font = title_font
        self.body_font = body_font
        self.title_size = Pt(title_size)
        self.body_size = Pt(body_size)
        self.accent_color = RGBColor(*accent_color)


THEMES = {
    "business_blue": ThemeConfig(
        primary_color=(0, 112, 192),
        secondary_color=(68, 84, 106),
        accent_color=(255, 192, 0)
    ),
    "business_dark": ThemeConfig(
        primary_color=(51, 51, 51),
        secondary_color=(102, 102, 102),
        background_color=(245, 245, 245),
        accent_color=(0, 176, 80)
    ),
    "education_green": ThemeConfig(
        primary_color=(0, 128, 0),
        secondary_color=(34, 139, 34),
        accent_color=(255, 215, 0)
    ),
    "creative_orange": ThemeConfig(
        primary_color=(255, 102, 0),
        secondary_color=(255, 153, 51),
        accent_color=(0, 102, 204)
    ),
    "tech_purple": ThemeConfig(
        primary_color=(102, 51, 153),
        secondary_color=(153, 102, 204),
        accent_color=(0, 204, 204)
    ),
    "minimal_white": ThemeConfig(
        primary_color=(0, 0, 0),
        secondary_color=(128, 128, 128),
        background_color=(255, 255, 255),
        accent_color=(200, 200, 200)
    )
}


class PPTGenerator:
    def __init__(self, theme_name: str = "business_blue"):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.theme = THEMES.get(theme_name, THEMES["business_blue"])
        self.image_service = None
        self.image_cache = {}
        self.image_config = {
            "enabled": True,
            "max_per_presentation": 5,
            "probability": 0.6,
            "retry_attempts": 2,
            "timeout": 120,
            "min_content_length": 20,
            "skip_keywords": ["总结", "目录", "大纲", "结束", "感谢", "联系方式", "谢谢"]
        }
    
    def set_theme(self, theme_config: ThemeConfig):
        self.theme = theme_config
    
    def _add_background(self, slide, color: Optional[RGBColor] = None):
        # slide.follow_master_background = False # 这一行在某些版本的 python-pptx 中不可用
        background = slide.background
        fill = background.fill
        fill.solid()
        # 强制设置背景色，并尝试处理颜色对象
        bg_color = color or self.theme.background_color
        try:
            fill.fore_color.rgb = bg_color
        except TypeError:
            # 如果直接赋值失败，尝试转换
            if isinstance(bg_color, tuple):
                fill.fore_color.rgb = RGBColor(*bg_color)
            else:
                # 最后的兜底
                print(f"[PPTGenerator] Warning: Could not set background color {bg_color}")
                
    def _add_decorative_shapes(self, slide, style: str = "default"):
        shapes_config = {
            "default": [
                {"type": MSO_SHAPE.OVAL, "left": 12.5, "top": 0.3, "width": 1.5, "height": 1.5},
                {"type": MSO_SHAPE.OVAL, "left": 11.8, "top": 6.0, "width": 2.0, "height": 2.0},
            ],
            "corner": [
                {"type": MSO_SHAPE.RIGHT_TRIANGLE, "left": 0, "top": 0, "width": 3, "height": 3},
                {"type": MSO_SHAPE.RIGHT_TRIANGLE, "left": 10.333, "top": 4.5, "width": 3, "height": 3},
            ],
            "circles": [
                {"type": MSO_SHAPE.OVAL, "left": -0.5, "top": 5.0, "width": 2.5, "height": 2.5},
                {"type": MSO_SHAPE.OVAL, "left": 11.0, "top": -0.5, "width": 3.0, "height": 3.0},
                {"type": MSO_SHAPE.OVAL, "left": 10.0, "top": 5.5, "width": 1.5, "height": 1.5},
            ],
            "bars": [
                {"type": MSO_SHAPE.RECTANGLE, "left": 0, "top": 0, "width": 0.3, "height": 7.5},
                {"type": MSO_SHAPE.RECTANGLE, "left": 12.8, "top": 0, "width": 0.5, "height": 7.5},
            ],
            "dots": [
                {"type": MSO_SHAPE.OVAL, "left": 0.3, "top": 0.3, "width": 0.3, "height": 0.3},
                {"type": MSO_SHAPE.OVAL, "left": 0.8, "top": 0.3, "width": 0.3, "height": 0.3},
                {"type": MSO_SHAPE.OVAL, "left": 1.3, "top": 0.3, "width": 0.3, "height": 0.3},
                {"type": MSO_SHAPE.OVAL, "left": 12.5, "top": 6.8, "width": 0.4, "height": 0.4},
                {"type": MSO_SHAPE.OVAL, "left": 12.0, "top": 6.9, "width": 0.3, "height": 0.3},
            ]
        }
        
        shapes = shapes_config.get(style, shapes_config["default"])
        
        for config in shapes:
            shape = slide.shapes.add_shape(
                config["type"],
                Inches(config["left"]),
                Inches(config["top"]),
                Inches(config["width"]),
                Inches(config["height"])
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.theme.primary_color
            shape.line.fill.background()
            shape.fill.fore_color.brightness = 0.6
    
    def _add_icon_shape(self, slide, icon_type: str, left: float, top: float, size: float = 0.6):
        icon_shapes = {
            "check": MSO_SHAPE.OVAL,
            "star": MSO_SHAPE.STAR_5_POINT,
            "arrow": MSO_SHAPE.RIGHT_ARROW,
            "diamond": MSO_SHAPE.DIAMOND,
            "hexagon": MSO_SHAPE.HEXAGON,
            "lightning": MSO_SHAPE.LIGHTNING_BOLT,
            "sun": MSO_SHAPE.SUN,
            "heart": MSO_SHAPE.HEART,
            "cloud": MSO_SHAPE.CLOUD,
            "gear": MSO_SHAPE.GEAR_6,
        }
        
        shape_type = icon_shapes.get(icon_type, MSO_SHAPE.OVAL)
        
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(left),
            Inches(top),
            Inches(size),
            Inches(size)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.theme.accent_color
        shape.line.fill.background()
        
        return shape
    
    def _add_shape_with_text(self, slide, left, top, width, height, text, font_size, font_color, font_name, bold=False, alignment=PP_ALIGN.LEFT, shape_type=MSO_SHAPE.RECTANGLE, fill_color=None, line_color=None):
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        if fill_color:
            shape.fill.solid()
            try:
                shape.fill.fore_color.rgb = fill_color
            except TypeError:
                if isinstance(fill_color, tuple):
                    shape.fill.fore_color.rgb = RGBColor(*fill_color)
        else:
            shape.fill.background()
        
        if line_color:
            try:
                shape.line.color.rgb = line_color
            except TypeError:
                if isinstance(line_color, tuple):
                    shape.line.color.rgb = RGBColor(*line_color)
        else:
            shape.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.auto_size = None
        
        # 确保文本存在
        if not text:
            text = " "
            
        p = tf.paragraphs[0]
        p.text = str(text) # 确保是字符串
        p.font.size = font_size
        try:
            p.font.color.rgb = font_color
        except TypeError:
             if isinstance(font_color, tuple):
                p.font.color.rgb = RGBColor(*font_color)
                
        p.font.name = font_name
        p.font.bold = bold
        p.alignment = alignment
        
        return shape
    
    def create_title_slide(self, title: str, subtitle: Optional[str] = None):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_background(slide)
        
        self._add_shape_with_text(
            slide,
            left=0.5,
            top=2.5,
            width=12.333,
            height=1.5,
            text=title,
            font_size=self.theme.title_size,
            font_color=self.theme.primary_color,
            font_name=self.theme.title_font,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        if subtitle:
            self._add_shape_with_text(
                slide,
                left=0.5,
                top=4.2,
                width=12.333,
                height=0.8,
                text=subtitle,
                font_size=Pt(24),
                font_color=self.theme.secondary_color,
                font_name=self.theme.body_font,
                alignment=PP_ALIGN.CENTER
            )
        
        return slide
    
    def create_content_slide(
        self,
        title: str,
        content: List[str],
        layout_type: str = "single_column",
        image_path: Optional[str] = None,
        charts: Optional[List[Dict[str, Any]]] = None,
        tables: Optional[List[Dict[str, Any]]] = None
    ):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_background(slide)
        
        self._add_shape_with_text(
            slide,
            left=0.5,
            top=0.3,
            width=12.333,
            height=0.8,
            text=title,
            font_size=Pt(36),
            font_color=self.theme.primary_color,
            font_name=self.theme.title_font,
            bold=True
        )
        
        # 预处理 content，移除 markdown 标记
        import re
        cleaned_content = []
        if content:
            for item in content:
                # 移除 **bold** 标记
                item = re.sub(r'\*\*(.*?)\*\*', r'\1', str(item))
                cleaned_content.append(item)
        content = cleaned_content
        
        # Layout logic:
        # Priority: Chart > Table > Image > Text Only
        
        if charts and len(charts) > 0:
            self._add_content_with_media(slide, content, media_type="chart", media_data=charts[0])
        elif tables and len(tables) > 0:
            self._add_content_with_media(slide, content, media_type="table", media_data=tables[0])
        elif image_path and os.path.exists(image_path):
            self._add_content_with_media(slide, content, media_type="image", media_data=image_path)
        elif layout_type == "single_column":
            self._add_single_column_content(slide, content)
        elif layout_type == "two_column":
            self._add_two_column_content(slide, content)
        elif layout_type == "three_column":
            self._add_three_column_content(slide, content)
        else:
            self._add_single_column_content(slide, content)
        
        return slide
    
    def _add_content_with_media(self, slide, content: List[str], media_type: str, media_data: Any):
        # Text Box (Left)
        text_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(7.0),
            Inches(5.5)
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"● {item}"
            p.font.size = self.theme.body_size
            p.font.color.rgb = self.theme.secondary_color
            p.font.name = self.theme.body_font
            p.space_after = Pt(12)
            p.level = 0
            
        # Media (Right)
        left = 8.0
        top = 1.5
        width = 4.8
        height = 4.0
        
        if media_type == "image":
            try:
                slide.shapes.add_picture(media_data, Inches(left), Inches(top), width=Inches(width))
            except Exception as e:
                print(f"[PPTGenerator] Failed to add image: {e}")
        elif media_type == "chart":
            try:
                self._add_chart(slide, media_data, left, top, width, height)
            except Exception as e:
                print(f"[PPTGenerator] Failed to add chart: {e}")
        elif media_type == "table":
            try:
                self._add_table(slide, media_data, left, top, width, height)
            except Exception as e:
                print(f"[PPTGenerator] Failed to add table: {e}")

    def _add_chart(self, slide, chart_data: Dict[str, Any], left: float, top: float, width: float, height: float):
        chart_type_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE
        }
        
        c_type = chart_type_map.get(chart_data.get("type", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
        c_data = CategoryChartData()
        
        data = chart_data.get("data", {})
        categories = data.get("categories", [])
        c_data.categories = categories
        
        series_list = data.get("series", [])
        for series in series_list:
            c_data.add_series(series.get("name", ""), series.get("values", []))
            
        chart = slide.shapes.add_chart(
            c_type, Inches(left), Inches(top), Inches(width), Inches(height), c_data
        ).chart
        
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        
        if chart_data.get("title"):
            chart.chart_title.text_frame.text = chart_data.get("title")

    def _add_table(self, slide, table_data: Dict[str, Any], left: float, top: float, width: float, height: float):
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        
        if not headers and not rows:
            return
            
        rows_count = len(rows) + 1
        cols_count = len(headers)
        
        shape = slide.shapes.add_table(
            rows_count, cols_count, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = shape.table
        
        # Set headers
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.primary_color
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
            
        # Set rows
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                if j < cols_count:
                    cell = table.cell(i + 1, j)
                    cell.text = str(val)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.CENTER

    def _add_content_with_image(self, slide, content: List[str], image_path: str):
        # This method is replaced by _add_content_with_media but kept for backward compatibility if needed
        # Or we can just remove it if we update all calls
        self._add_content_with_media(slide, content, "image", image_path)

    def _add_single_column_content(self, slide, content: List[str]):
        text_box = slide.shapes.add_textbox(
            Inches(0.5),
            Inches(1.5),
            Inches(12.333),
            Inches(5.5)
        )
        
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(content):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"● {item}"
            p.font.size = self.theme.body_size
            p.font.color.rgb = self.theme.secondary_color
            p.font.name = self.theme.body_font
            p.space_after = Pt(12)
            p.level = 0
    
    def _add_two_column_content(self, slide, content: List[str]):
        mid = len(content) // 2
        left_content = content[:mid]
        right_content = content[mid:]
        
        for col, (x_offset, items) in enumerate([(0.5, left_content), (6.917, right_content)]):
            text_box = slide.shapes.add_textbox(
                Inches(x_offset),
                Inches(1.5),
                Inches(5.917),
                Inches(5.5)
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = self.theme.body_size
                p.font.color.rgb = self.theme.secondary_color
                p.font.name = self.theme.body_font
                p.space_after = Pt(12)
    
    def _add_three_column_content(self, slide, content: List[str]):
        third = len(content) // 3
        columns = [
            content[:third],
            content[third:2*third],
            content[2*third:]
        ]
        
        x_positions = [0.5, 4.611, 8.722]
        
        for x_offset, items in zip(x_positions, columns):
            text_box = slide.shapes.add_textbox(
                Inches(x_offset),
                Inches(1.5),
                Inches(3.722),
                Inches(5.5)
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = f"• {item}"
                p.font.size = Pt(16)
                p.font.color.rgb = self.theme.secondary_color
                p.font.name = self.theme.body_font
                p.space_after = Pt(10)
    
    def create_section_slide(self, title: str, section_number: int):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_background(slide)
        
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.917),
            Inches(2.0),
            Inches(1.5),
            Inches(1.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = self.theme.primary_color
        num_shape.line.fill.background()
        
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        p.text = str(section_number)
        p.font.size = Pt(48)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.anchor = MSO_ANCHOR.MIDDLE
        
        self._add_shape_with_text(
            slide,
            left=0.5,
            top=4.0,
            width=12.333,
            height=1.0,
            text=title,
            font_size=Pt(36),
            font_color=self.theme.primary_color,
            font_name=self.theme.title_font,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        return slide
    
    def create_ending_slide(self, title: str = "感谢观看", contact_info: Optional[str] = None):
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._add_background(slide)
        
        self._add_shape_with_text(
            slide,
            left=0.5,
            top=2.5,
            width=12.333,
            height=1.5,
            text=title,
            font_size=self.theme.title_size,
            font_color=self.theme.primary_color,
            font_name=self.theme.title_font,
            bold=True,
            alignment=PP_ALIGN.CENTER
        )
        
        if contact_info:
            self._add_shape_with_text(
                slide,
                left=0.5,
                top=4.5,
                width=12.333,
                height=0.8,
                text=contact_info,
                font_size=Pt(20),
                font_color=self.theme.secondary_color,
                font_name=self.theme.body_font,
                alignment=PP_ALIGN.CENTER
            )
        
        return slide
        
        return slide
    
    def generate_from_json(self, content_json: Dict[str, Any], generate_images: bool = True, progress_callback=None) -> Presentation:
        title = content_json.get("title", "演示文稿")
        pages = content_json.get("pages", [])
        
        theme_settings = content_json.get("theme")
        
        if theme_settings:
            if isinstance(theme_settings, dict):
                try:
                    primary_color = theme_settings.get("primary_color", (0, 112, 192))
                    secondary_color = theme_settings.get("secondary_color", (68, 84, 106))
                    background_color = theme_settings.get("background_color", (255, 255, 255))
                    accent_color = theme_settings.get("accent_color", (255, 192, 0))
                    
                    if isinstance(primary_color, list):
                        primary_color = tuple(primary_color)
                    if isinstance(secondary_color, list):
                        secondary_color = tuple(secondary_color)
                    if isinstance(background_color, list):
                        background_color = tuple(background_color)
                    if isinstance(accent_color, list):
                        accent_color = tuple(accent_color)
                    
                    custom_theme = ThemeConfig(
                        primary_color=primary_color,
                        secondary_color=secondary_color,
                        background_color=background_color,
                        title_font=theme_settings.get("title_font", "微软雅黑"),
                        body_font=theme_settings.get("body_font", "微软雅黑"),
                        title_size=theme_settings.get("title_size", 44),
                        body_size=theme_settings.get("body_size", 18),
                        accent_color=accent_color
                    )
                    self.set_theme(custom_theme)
                    logger.info(f"Applied custom theme from JSON: {theme_settings.get('theme', 'custom')}")
                except Exception as e:
                    logger.error(f"Failed to apply custom theme: {e}")
            elif isinstance(theme_settings, str) and theme_settings in THEMES:
                self.theme = THEMES[theme_settings]
                logger.info(f"Applied preset theme: {theme_settings}")
        
        if generate_images and self.image_config["enabled"]:
            self.image_service = ImageServiceFactory.get_service()
            logger.info(f"Image service initialized: {type(self.image_service).__name__}")
        
        if progress_callback:
            progress_callback(50, "正在创建标题页...")
        
        self.create_title_slide(title)
        
        generated_count = 0
        total_pages = len(pages)
        
        for i, page in enumerate(pages):
            page_progress = 50 + int((i / total_pages) * 35) if total_pages > 0 else 50
            page_title = page.get("title", f"第{i+1}页")
            
            if progress_callback:
                progress_callback(page_progress, f"正在生成第 {i+1}/{total_pages} 页: {page_title}...")
            
            page_content = page.get("content", [])
            layout_type = page.get("layout_type", "single_column")
            
            charts = page.get("charts", [])
            tables = page.get("tables", [])
            
            image_path = None
            
            should_generate = self._should_generate_image_for_page(
                page_title=page_title,
                page_content=page_content,
                page_index=i,
                total_pages=total_pages,
                has_charts=bool(charts),
                has_tables=bool(tables),
                generated_count=generated_count,
                generate_images=generate_images
            )
            
            logger.info(f"Page {i+1}/{total_pages}: should_generate={should_generate}, has_charts={bool(charts)}, has_tables={bool(tables)}, generated_count={generated_count}")
            
            if should_generate and self.image_service:
                image_prompt = self._generate_image_prompt(title, page_title, page_content)
                logger.info(f"Generating image with prompt: {image_prompt[:100]}...")
                
                image_path = self._generate_image_with_retry(image_prompt)
                
                if image_path:
                    generated_count += 1
                    logger.info(f"Successfully generated image for page {i+1}, path: {image_path}")
                    page["image_path"] = image_path
            
            self.create_content_slide(page_title, page_content, layout_type, page.get("image_path"), charts, tables)
        
        if progress_callback:
            progress_callback(85, "正在创建结束页...")
        
        self.create_ending_slide()
        
        logger.info(f"PPT generation completed. Total images generated: {generated_count}")
        
        return self.prs
    
    def _should_generate_image_for_page(
        self,
        page_title: str,
        page_content: List[str],
        page_index: int,
        total_pages: int,
        has_charts: bool,
        has_tables: bool,
        generated_count: int,
        generate_images: bool
    ) -> bool:
        logger.info(f"_should_generate_image_for_page called: generate_images={generate_images}, enabled={self.image_config['enabled']}")
        
        if not generate_images or not self.image_config["enabled"]:
            logger.info(f"Skipping image generation - disabled (generate_images={generate_images}, enabled={self.image_config['enabled']})")
            return False
        
        if has_charts or has_tables:
            logger.info(f"Skipping image generation - page has charts or tables")
            return False
        
        if generated_count >= self.image_config["max_per_presentation"]:
            logger.info(f"Skipping image generation - max limit reached ({self.image_config['max_per_presentation']})")
            return False
        
        skip_keywords = self.image_config["skip_keywords"]
        for keyword in skip_keywords:
            if keyword in page_title:
                logger.info(f"Skipping image generation - contains keyword: {keyword}")
                return False
        
        content_text = " ".join(page_content) if page_content else ""
        if len(content_text) < self.image_config["min_content_length"]:
            logger.info(f"Skipping image generation - content too short ({len(content_text)} chars)")
            return False
        
        random_val = random.random()
        logger.info(f"Random value: {random_val}, threshold: {self.image_config['probability']}")
        
        if random_val >= self.image_config["probability"]:
            logger.info(f"Skipping image generation - probability check failed")
            return False
        
        logger.info(f"Will generate image for page: {page_title}")
        return True
    
    def _generate_image_with_retry(self, prompt: str) -> Optional[str]:
        cache_key = hashlib.md5(prompt.encode()).hexdigest()
        
        if cache_key in self.image_cache:
            logger.info(f"Using cached image for prompt")
            return self.image_cache[cache_key]
        
        for attempt in range(self.image_config["retry_attempts"]):
            try:
                import threading
                result_container = {'result': None, 'error': None}
                
                def run_async_in_thread():
                    try:
                        new_loop = asyncio.new_event_loop()
                        asyncio.set_event_loop(new_loop)
                        try:
                            result_container['result'] = new_loop.run_until_complete(
                                self.image_service.generate_image(prompt)
                            )
                        finally:
                            new_loop.close()
                    except Exception as e:
                        result_container['error'] = e
                
                thread = threading.Thread(target=run_async_in_thread)
                thread.start()
                thread.join(timeout=self.image_config["timeout"])
                
                if result_container['error']:
                    raise result_container['error']
                
                image_path = result_container['result']
                
                if image_path:
                    self.image_cache[cache_key] = image_path
                    logger.info(f"Image generated successfully on attempt {attempt + 1}")
                    return image_path
                else:
                    logger.warning(f"Image generation returned None on attempt {attempt + 1}")
                    
            except Exception as e:
                logger.error(f"Image generation failed on attempt {attempt + 1}: {e}")
                if attempt < self.image_config["retry_attempts"] - 1:
                    import time
                    time.sleep(2)
                    continue
        
        logger.error(f"Image generation failed after {self.image_config['retry_attempts']} attempts")
        return None
    
    def _should_generate_image(self, page_title: str, page_index: int, total_pages: int) -> bool:
        return random.random() < 0.5
    
    def _generate_image_prompt(self, ppt_title: str, page_title: str, content: List[str]) -> str:
        content_text = " ".join(content) if content else ""
        
        style_keywords = {
            "科技": "technology, futuristic, digital, innovation, circuit patterns, blue neon lights",
            "商务": "business, corporate, professional, office, teamwork, modern workspace",
            "教育": "education, learning, books, knowledge, academic, classroom",
            "医疗": "medical, healthcare, hospital, medicine, doctor, health",
            "金融": "finance, banking, money, investment, stock market, wealth",
            "营销": "marketing, advertising, social media, digital marketing, growth",
            "环保": "environment, green energy, sustainability, nature, eco-friendly",
            "创新": "innovation, creativity, idea, lightbulb, breakthrough, startup"
        }
        
        detected_style = None
        for keyword, style in style_keywords.items():
            if keyword in ppt_title or keyword in page_title or keyword in content_text:
                detected_style = style
                break
        
        prompt_parts = []
        
        prompt_parts.append("Professional presentation background image")
        
        if detected_style:
            prompt_parts.append(f"with {detected_style} theme")
        
        prompt_parts.append(f"related to: {page_title}")
        
        if content_text:
            keywords = self._extract_keywords(content_text, max_keywords=3)
            if keywords:
                prompt_parts.append(f"keywords: {', '.join(keywords)}")
        
        prompt_parts.extend([
            "Style: Modern, Clean, Professional",
            "Quality: High resolution, 4K, photorealistic",
            "Requirements: No text, no watermark, no people faces, abstract or conceptual",
            "Color scheme: Matching business presentation standards"
        ])
        
        prompt = ". ".join(prompt_parts) + "."
        
        return prompt
    
    def _extract_keywords(self, text: str, max_keywords: int = 3) -> List[str]:
        import re
        
        stop_words = {
            "的", "是", "在", "和", "了", "有", "我", "他", "她", "它", "们",
            "这", "那", "就", "也", "都", "而", "及", "与", "或", "但",
            "可以", "需要", "应该", "能够", "通过", "进行", "实现", "使用",
            "一个", "这个", "那个", "这些", "那些", "其", "之", "以",
            "为", "于", "上", "下", "中", "来", "去", "到", "从", "向"
        }
        
        words = re.findall(r'[\u4e00-\u9fa5]{2,4}|[a-zA-Z]{3,}', text)
        
        word_freq = {}
        for word in words:
            if word.lower() not in stop_words and word not in stop_words:
                word_freq[word] = word_freq.get(word, 0) + 1
        
        sorted_words = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)
        
        keywords = [word for word, freq in sorted_words[:max_keywords]]
        
        return keywords
    
    def save(self, file_path: str):
        self.prs.save(file_path)
    
    def save_to_bytes(self) -> bytes:
        buffer = io.BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()
    
    def get_presentation(self) -> Presentation:
        return self.prs


def generate_ppt(content_json: Dict[str, Any], theme: str = "business_blue") -> bytes:
    generator = PPTGenerator(theme)
    generator.generate_from_json(content_json)
    return generator.save_to_bytes()
